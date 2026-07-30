"""
Cortex AES PoV wrap-up deck.

Ported from jbillochon/povplatform (rendering/deck.py): same palette, same
20 x 11.25in 16:9 geometry, same native editable shapes - no rasterised
images. Added here: a success-criteria scorecard, threat-intelligence and
XSIAM slides, plus null-awareness so an uncollected domain reads
"not measured" rather than zero.

The analytical slides (supply chain, findings, attack scenarios, recommended
actions, threat context) live in deck_sections.py and are bound onto
DeckBuilder by its attach() at the bottom of this module; build() calls them
in place so the slide order stays readable here.
"""

from __future__ import annotations

import logging
from datetime import date as _date

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from . import deck_sections
from .common import NOT_MEASURED, Data

log = logging.getLogger(__name__)

# ---------------------------------------------------------------- palette
BG = RGBColor(0x04, 0x16, 0x10)
BG_CARD = RGBColor(0x0C, 0x27, 0x1C)
BG_CARD2 = RGBColor(0x13, 0x35, 0x27)
BG_MUTED = RGBColor(0x0A, 0x1F, 0x17)
BORDER = RGBColor(0x1F, 0x4A, 0x36)
MINT = RGBColor(0x00, 0xD2, 0x6A)
TEAL = RGBColor(0x5F, 0xE3, 0xC0)
AMBER = RGBColor(0xF5, 0xB0, 0x17)
RED = RGBColor(0xFF, 0x5C, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MID = RGBColor(0xB4, 0xCF, 0xC0)
TEXT_LO = RGBColor(0x74, 0x8F, 0x7F)

FONT = "Calibri"
MONO = "Consolas"

RISK_COLORS = {
    "critical": RED,
    "high": AMBER,
    "medium": TEAL,
    "low": MINT,
    "pending": TEXT_LO,
}

VERDICT_COLORS = {
    "met": MINT,
    "passed": MINT,
    "partially met": AMBER,
    "partial": AMBER,
    "not met": RED,
    "failed": RED,
    "not tested": TEXT_LO,
    "not measured": TEXT_LO,
}

# Slide geometry (inches)
SW, SH = 20.0, 11.25
M = 0.75
CONTENT_W = SW - 2 * M


def _n(value) -> str:
    try:
        return format(int(value or 0), ",")
    except (TypeError, ValueError):
        return str(value)


class DeckBuilder:
    def __init__(self, data: dict, narrative: dict | None = None):
        self.r = Data(data)
        self.n = narrative or {}
        self.prs = Presentation()
        self.prs.slide_width = Inches(SW)
        self.prs.slide_height = Inches(SH)
        self.blank = self.prs.slide_layouts[6]

    # ------------------------------------------------------------ helpers
    @staticmethod
    def _no_shadow(shape):
        try:
            spPr = shape._element.spPr
            for el in spPr.findall(qn("a:effectLst")):
                spPr.remove(el)
            etree.SubElement(spPr, qn("a:effectLst"))
        except Exception:  # noqa: BLE001 - cosmetic only
            pass

    def _slide(self):
        s = self.prs.slides.add_slide(self.blank)
        bg = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG
        bg.line.fill.background()
        self._no_shadow(bg)
        return s

    def _card(self, s, x, y, w, h, fill=BG_CARD, border=BORDER, radius=0.04, bw=1.0):
        sh = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if border:
            sh.line.color.rgb = border
            sh.line.width = Pt(bw)
        else:
            sh.line.fill.background()
        sh.adjustments[0] = radius
        self._no_shadow(sh)
        return sh

    def _text(self, s, x, y, w, h, text, size=14, color=WHITE, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
              font=FONT, line_spacing=None):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, line in enumerate(str(text).split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if line_spacing:
                p.line_spacing = line_spacing
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
        return tb

    def _pill(self, s, x, y, w, h, text, fill=MINT, color=BG, size=11):
        p = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        p.fill.solid()
        p.fill.fore_color.rgb = fill
        p.line.fill.background()
        p.adjustments[0] = 0.5
        self._no_shadow(p)
        tf = p.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        r = para.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = FONT
        return p

    def _header(self, s, title, subtitle="", tag="POV RESULTS", tag_color=MINT):
        self._pill(s, M, 0.55, 1.9, 0.34, tag, tag_color, BG, 11)
        self._text(s, M, 1.00, CONTENT_W, 0.62, title, size=32, bold=True)
        if subtitle:
            self._text(s, M, 1.66, CONTENT_W, 0.36, subtitle, size=14, color=TEXT_MID)

    def _footer(self, s, note=""):
        left = "Cortex AES  \u00b7  " + self.r.customer + "  \u00b7  PoV wrap-up"
        self._text(s, M, SH - 0.55, CONTENT_W * 0.6, 0.28, left, size=9, color=TEXT_LO)
        if note:
            self._text(s, M + CONTENT_W * 0.6, SH - 0.55, CONTENT_W * 0.4, 0.28,
                       note, size=9, color=TEXT_LO, align=PP_ALIGN.RIGHT, italic=True)

    def _kpi(self, s, x, y, w, h, value, label, sub="", accent=MINT):
        self._card(s, x, y, w, h, fill=BG_CARD)
        size = 54 if len(str(value)) <= 7 else 28
        self._text(s, x + 0.4, y + 0.35, w - 0.8, 1.1, str(value),
                   size=size, bold=True, color=accent)
        self._text(s, x + 0.4, y + 1.55, w - 0.8, 0.42, label, size=16, bold=True)
        if sub:
            self._text(s, x + 0.4, y + 1.98, w - 0.8, 0.4, sub,
                       size=13, color=TEXT_MID, italic=True)

    def _bar_chart(self, s, x, y, w, h, data, accent=MINT,
                   title="", max_rows=8, color_map=None):
        """Horizontal bar chart drawn with native rectangles."""
        if title:
            self._text(s, x, y, w, 0.38, title, size=17, bold=True, color=TEXT_MID)
            y += 0.60
            h -= 0.60
        if data is None:
            self._text(s, x, y, w, 0.4, NOT_MEASURED + " during this PoV",
                       size=12, color=TEXT_LO, italic=True)
            return
        rows = list(data.items())[:max_rows]
        if not rows:
            self._text(s, x, y, w, 0.4, "No data collected", size=12,
                       color=TEXT_LO, italic=True)
            return
        vmax = max(v for _, v in rows) or 1
        row_h = min(0.72, h / max(len(rows), 1))
        gap = row_h * 0.24
        label_w = w * 0.34
        bar_zone = w - label_w - 1.1

        for i, (k, v) in enumerate(rows):
            ry = y + i * row_h
            self._text(s, x, ry, label_w - 0.15, row_h - gap, str(k)[:38],
                       size=13, color=TEXT_MID, anchor=MSO_ANCHOR.MIDDLE)
            bw = max(0.06, bar_zone * (v / vmax))
            col = (color_map or {}).get(str(k).lower(), accent)
            bar = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + label_w), Inches(ry + gap * 0.35),
                Inches(bw), Inches(row_h - gap),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = col
            bar.line.fill.background()
            bar.adjustments[0] = 0.35
            self._no_shadow(bar)
            self._text(s, x + label_w + bw + 0.15, ry, 1.3, row_h - gap,
                       _n(v), size=13, bold=True, color=WHITE,
                       anchor=MSO_ANCHOR.MIDDLE)

    def _table(self, s, x, y, w, headers, rows, col_w, row_h=0.58,
               max_rows=10, risk_col=None, color_map=None):
        hy = y
        self._card(s, x, hy, w, row_h, fill=BG_CARD2, border=BORDER, radius=0.06)
        cx = x + 0.18
        for i, htxt in enumerate(headers):
            self._text(s, cx, hy + 0.02, col_w[i], row_h - 0.04, htxt,
                       size=13, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[i]
        for ri, row in enumerate(rows[:max_rows]):
            ry = hy + row_h + ri * row_h
            if ri % 2 == 0:
                self._card(s, x, ry, w, row_h, fill=BG_MUTED, border=None, radius=0.02)
            cx = x + 0.18
            for ci, cell in enumerate(row):
                color = WHITE
                bold = ci == 0
                if risk_col is not None and ci == risk_col:
                    palette = color_map or RISK_COLORS
                    color = palette.get(str(cell).lower(), WHITE)
                    bold = True
                self._text(s, cx, ry + 0.02, col_w[ci], row_h - 0.04, str(cell),
                           size=12.5, color=color, bold=bold, anchor=MSO_ANCHOR.MIDDLE)
                cx += col_w[ci]

    def _empty_card(self, s, x, y, w, text):
        self._card(s, x, y, w, 1.4)
        self._text(s, x + 0.5, y + 0.5, w - 1.0, 0.5, text,
                   size=14, color=TEXT_LO, italic=True)

    # ------------------------------------------------------------- slides
    def slide_title(self):
        s = self._slide()
        self._pill(s, M, 3.6, 2.4, 0.42, "POV WRAP-UP", MINT, BG, 13)
        self._text(s, M, 4.2, CONTENT_W, 1.2, "Cortex AES", size=64, bold=True)
        self._text(s, M, 5.4, CONTENT_W, 0.8,
                   "Proof of Value Results \u2014 " + self.r.customer,
                   size=28, color=TEAL)
        window = ""
        start = self.r.meta.get("pov_start")
        end = self.r.meta.get("pov_end")
        if start and end:
            window = str(start) + " \u2192 " + str(end)
        prepared = self.r.meta.get("prepared_by") or ""
        line = "  \u00b7  ".join(x for x in [window, prepared] if x)
        if line:
            self._text(s, M, 6.3, CONTENT_W, 0.4, line, size=14, color=TEXT_MID)
        generated = self.r.meta.get("generated_at") or ""
        stamp = generated[:10] if generated else _date.today().isoformat()
        self._text(s, M, SH - 1.1, CONTENT_W, 0.4,
                   "Generated " + stamp + " from live Koi tenant data",
                   size=11, color=TEXT_LO, italic=True)
        return s

    def slide_exec_summary(self):
        s = self._slide()
        self._header(s, "Executive Summary",
                     "What the platform found across your fleet during the PoV")
        y = 2.5
        cw = (CONTENT_W - 3 * 0.35) / 4
        kpi_h = 2.7
        crit = self.r.critical
        remediated = self.r.mapping("remediations_by_status").get("remediated", 0)

        self._kpi(s, M, y, cw, kpi_h, self.r.num("devices_total"),
                  "Endpoints covered",
                  self.r.num("devices_active") + " active", MINT)
        self._kpi(s, M + (cw + 0.35), y, cw, kpi_h, self.r.num("items_total"),
                  "Items discovered",
                  self.r.num("unique_publishers") + " unique publishers", TEAL)
        self._kpi(s, M + 2 * (cw + 0.35), y, cw, kpi_h, self.r.high_and_critical,
                  "High & critical risk items",
                  (_n(crit) + " critical") if self.r.measured("items_by_risk") else "",
                  RED if crit else AMBER)
        remediated_txt = (_n(remediated) if self.r.measured("remediations_total")
                          else NOT_MEASURED)
        self._kpi(s, M + 3 * (cw + 0.35), y, cw, kpi_h, remediated_txt,
                  "Items remediated",
                  self.r.num("remediations_total") + " flagged total", MINT)

        y2 = y + kpi_h + 0.45
        h2 = SH - y2 - 0.9
        cw2 = (CONTENT_W - 2 * 0.35) / 3
        pad = 0.4

        self._card(s, M, y2, cw2, h2)
        self._bar_chart(s, M + pad, y2 + 0.5, cw2 - 2 * pad, h2 - 1.0,
                        self.r.get("items_by_view"), TEAL,
                        title="Attack surface", max_rows=7)

        x3 = M + cw2 + 0.35
        self._card(s, x3, y2, cw2, h2)
        self._bar_chart(s, x3 + pad, y2 + 0.5, cw2 - 2 * pad, h2 - 1.0,
                        self.r.get("items_by_risk"), AMBER,
                        title="Risk distribution", max_rows=5,
                        color_map=RISK_COLORS)

        x4 = M + 2 * (cw2 + 0.35)
        self._card(s, x4, y2, cw2, h2)
        gx = x4 + pad
        self._text(s, gx, y2 + 0.5, cw2 - 2 * pad, 0.4,
                   "Governance in place", size=15, bold=True, color=TEXT_MID)
        policies = self.r.rows("policies")
        ungoverned = self.r.get("ungoverned_high_risk")
        gov_lines = [
            (self.r.num("policies_enabled") + " of " + str(len(policies))
             + " marketplace policies enabled", TEXT_MID),
            (str(len(self.r.rows("runtime_policies")))
             + " agent runtime policies active", TEXT_MID),
            (self.r.num("allowlist_count") + " items allowlisted", TEXT_MID),
            (self.r.num("blocklist_count") + " items blocklisted", TEXT_MID),
            (self.r.num("agent_sessions_total") + " agent sessions observed", TEXT_MID),
            (self.r.num("ungoverned_high_risk") + " high-risk items still ungoverned",
             AMBER if ungoverned else TEXT_MID),
        ]
        gy = y2 + 1.15
        step = min(0.62, (h2 - 1.6) / max(len(gov_lines), 1))
        for i, (line, col) in enumerate(gov_lines):
            self._text(s, gx, gy + i * step, cw2 - 2 * pad, step,
                       "\u25b8  " + line, size=13, color=col,
                       anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s)
        return s

    def slide_narrative(self):
        """Operator-written executive summary, when there is one."""
        text = self.n.get("executive_summary")
        if not text:
            return None
        s = self._slide()
        self._header(s, "In summary", "Written from the collected tenant data")
        self._card(s, M, 2.5, CONTENT_W, SH - 3.4)
        self._text(s, M + 0.6, 3.0, CONTENT_W - 1.2, SH - 4.2, text,
                   size=17, color=TEXT_MID, line_spacing=1.35)
        self._footer(s)
        return s

    def slide_scorecard(self):
        """Success criteria agreed at kickoff. The heart of the deliverable."""
        criteria = self.n.get("success_criteria") or []
        s = self._slide()
        self._header(s, "Success criteria",
                     "What we agreed to prove at kickoff, and what the data shows",
                     tag="SCORECARD", tag_color=TEAL)
        if not criteria:
            self._empty_card(
                s, M, 2.7, CONTENT_W,
                "[[TO BE PROVIDED: success criteria agreed at kickoff]] \u2014 "
                "none were supplied, so no verdict can be shown here.",
            )
            self._footer(s)
            return s
        rows = [
            [str(c.get("criterion", ""))[:78],
             str(c.get("verdict", "")).capitalize(),
             str(c.get("evidence", ""))[:70]]
            for c in criteria[:9]
        ]
        self._table(s, M, 2.7, CONTENT_W,
                    ["Criterion", "Verdict", "Evidence"], rows,
                    col_w=[8.0, 2.4, 8.1], row_h=0.72, max_rows=9,
                    risk_col=1, color_map=VERDICT_COLORS)
        self._footer(s, "Every verdict traces to collected tenant data")
        return s

    def slide_discovery(self):
        s = self._slide()
        self._header(s, "Discovery \u2014 what is actually installed",
                     "Full inventory across every governed marketplace and registry")
        y = 2.5
        h = SH - y - 0.9
        half = (CONTENT_W - 0.4) / 2
        self._card(s, M, y, half, h)
        self._bar_chart(s, M + 0.5, y + 0.5, half - 1.0, h - 1.0,
                        self.r.get("items_by_view"), TEAL,
                        title="Items by category", max_rows=9)
        self._card(s, M + half + 0.4, y, half, h)
        self._bar_chart(s, M + half + 0.9, y + 0.5, half - 1.0, h - 1.0,
                        self.r.get("items_by_marketplace"), MINT,
                        title="Items by marketplace", max_rows=10)
        self._footer(s, "Source: Koi inventory API")
        return s

    def slide_risk(self):
        s = self._slide()
        self._header(
            s, "Risk \u2014 where the exposure is",
            self.r.high_and_critical
            + " items rated high or critical by the risk engine",
            tag="RISK FINDINGS", tag_color=AMBER,
        )
        y = 2.5
        rows = []
        for t in self.r.rows("top_risk_items"):
            score = t.get("risk")
            score_txt = format(score, ".1f") if isinstance(score, (int, float)) else "-"
            rows.append([
                (t.get("name") or "")[:42],
                (t.get("publisher") or "")[:22],
                t.get("marketplace") or "",
                (t.get("risk_level") or "").capitalize(),
                score_txt,
                _n(t.get("endpoints")),
            ])
        tw = CONTENT_W * 0.60
        self._text(s, M, y, tw, 0.4, "Highest-risk items discovered",
                   size=17, bold=True, color=TEXT_MID)
        if rows:
            self._table(
                s, M, y + 0.55, tw,
                ["Item", "Publisher", "Marketplace", "Risk", "Score", "Endpoints"],
                rows, col_w=[3.7, 2.1, 2.3, 1.3, 1.1, 1.4],
                max_rows=10, risk_col=3,
            )
        else:
            self._empty_card(s, M, y + 0.55, tw,
                             "Inventory was not collected during this PoV")
        fx = M + tw + 0.5
        fw = CONTENT_W - tw - 0.5
        fh = SH - (y + 0.55) - 0.9
        self._card(s, fx, y + 0.55, fw, fh)
        freq = self.r.get("finding_frequency")
        freq_map = None
        if freq is not None:
            freq_map = {f.get("finding"): f.get("count") for f in freq}
        self._bar_chart(s, fx + 0.45, y + 1.05, fw - 0.9, fh - 1.0, freq_map,
                        AMBER, title="Most frequent findings", max_rows=8)
        self._footer(s, "Risk scores from the Koi risk engine (0-10)")
        return s

    def slide_threat_intel(self):
        """Deterministic TI: KEV first, then EPSS, then CVSS."""
        enr = self.r.enrichment
        if not enr:
            return None
        s = self._slide()
        fetched = (enr.get("fetched_at") or "")[:10]
        subtitle = "Exploitation first: KEV, then EPSS probability, then CVSS severity"
        if fetched:
            subtitle += "  \u00b7  as of " + fetched
        self._header(s, "Threat intelligence", subtitle,
                     tag="THREAT INTEL", tag_color=RED)
        y = 2.5
        cw = (CONTENT_W - 3 * 0.35) / 4
        kpi_h = 2.3
        cves = enr.get("cves") or {}
        osv = (enr.get("osv") or {}).get("matches") or {}
        kev = self.r.kev_count
        self._kpi(s, M, y, cw, kpi_h, _n(len(cves)), "CVEs in scope",
                  "from findings and package matches", TEAL)
        self._kpi(s, M + cw + 0.35, y, cw, kpi_h, _n(kev),
                  "Known exploited (KEV)",
                  "CISA catalogue" if kev else "none evidenced as exploited",
                  RED if kev else MINT)
        self._kpi(s, M + 2 * (cw + 0.35), y, cw, kpi_h, _n(len(osv)),
                  "Vulnerable packages", "exact version match (OSV)", AMBER)
        self._kpi(s, M + 3 * (cw + 0.35), y, cw, kpi_h,
                  _n(len(enr.get("mitre") or {})),
                  "Findings mapped to ATT&CK", "curated mapping", MINT)

        y2 = y + kpi_h + 0.45
        half = (CONTENT_W - 0.5) / 2
        self._text(s, M, y2, half, 0.4, "CVEs by exploitation evidence",
                   size=17, bold=True, color=TEXT_MID)
        ti = self.r.ti_rows(9)
        if ti:
            self._table(s, M, y2 + 0.55, half,
                        ["CVE", "KEV", "EPSS", "CVSS", "Severity"], ti,
                        col_w=[2.9, 1.2, 1.5, 1.2, 2.4], max_rows=9,
                        risk_col=4)
        else:
            self._empty_card(s, M, y2 + 0.55, half,
                             "No CVE identified in the collected findings")

        bx = M + half + 0.5
        self._text(s, bx, y2, half, 0.4, "ATT&CK techniques from findings",
                   size=17, bold=True, color=TEXT_MID)
        mitre = [[r[0][:38], r[1][:24], r[2][:34]] for r in self.r.mitre_rows(9)]
        if mitre:
            self._table(s, bx, y2 + 0.55, half,
                        ["Finding", "Technique", "Name"], mitre,
                        col_w=[3.6, 2.2, 3.4], max_rows=9)
        else:
            self._empty_card(s, bx, y2 + 0.55, half,
                             "No finding matched the curated ATT&CK mapping")
        self._footer(s, "Sources: NVD \u00b7 OSV.dev \u00b7 CISA KEV \u00b7 FIRST EPSS")
        return s

    def slide_governance(self):
        s = self._slide()
        self._header(s, "Governance \u2014 the controls you put in place",
                     "Policies, guardrails and lists configured during the PoV")
        y = 2.5
        half = (CONTENT_W - 0.5) / 2
        self._text(s, M, y, half, 0.4, "Marketplace policies", size=17,
                   bold=True, color=TEXT_MID)
        prows = []
        for p in self.r.rows("policies")[:10]:
            groups = p.get("groups")
            prows.append([
                (p.get("name") or "")[:38],
                (p.get("action") or "").capitalize(),
                "Enabled" if p.get("enabled") else "Disabled",
                str(groups) if groups else "All",
            ])
        if prows:
            self._table(s, M, y + 0.55, half,
                        ["Policy", "Action", "State", "Groups"],
                        prows, col_w=[4.5, 1.6, 1.6, 1.3], max_rows=8)
        else:
            self._empty_card(s, M, y + 0.55, half,
                             "No marketplace policy configured during this PoV")

        rx = M + half + 0.5
        self._text(s, rx, y, half, 0.4, "Agent runtime policies", size=17,
                   bold=True, color=TEXT_MID)
        rrows = []
        for p in self.r.rows("runtime_policies")[:10]:
            agents = ", ".join(p.get("agents") or [])[:24] or "-"
            kinds = ", ".join(sorted(set(p.get("rule_types") or [])))[:22] or "-"
            rrows.append([(p.get("name") or "")[:34],
                          (p.get("mode") or "").capitalize(), agents, kinds])
        if rrows:
            self._table(s, rx, y + 0.55, half,
                        ["Policy", "Mode", "Agents", "Rule types"],
                        rrows, col_w=[3.4, 1.4, 2.4, 1.8], max_rows=8)
        else:
            self._empty_card(s, rx, y + 0.55, half,
                             "No agent runtime policies configured during this PoV")

        ly = SH - 3.1
        cw = (CONTENT_W - 2 * 0.35) / 3
        approvals = self.r.mapping("approvals_by_status")
        self._kpi(s, M, ly, cw, 2.4, self.r.num("allowlist_count"),
                  "Items allowlisted", "explicitly approved", MINT)
        self._kpi(s, M + cw + 0.35, ly, cw, 2.4, self.r.num("blocklist_count"),
                  "Items blocklisted", "explicitly denied", RED)
        self._kpi(s, M + 2 * (cw + 0.35), ly, cw, 2.4,
                  _n(approvals.get("approved", 0)),
                  "Approval requests handled",
                  _n(approvals.get("pending", 0)) + " still pending", TEAL)
        self._footer(s)
        return s

    def slide_remediation(self):
        s = self._slide()
        self._header(s, "Remediation \u2014 risk actually removed",
                     "Items flagged and cleaned from endpoints during the PoV",
                     tag="OUTCOMES", tag_color=MINT)
        y = 2.5
        h = SH - y - 0.9
        left = CONTENT_W * 0.36
        self._card(s, M, y, left, h)
        self._bar_chart(s, M + 0.5, y + 0.5, left - 1.0, (h - 1.4) / 2,
                        self.r.get("remediations_by_status"), MINT,
                        title="Remediation status", max_rows=5)
        self._bar_chart(s, M + 0.5, y + 0.9 + (h - 1.4) / 2, left - 1.0,
                        (h - 1.4) / 2, self.r.get("alerts_by_severity"), AMBER,
                        title="Alerts raised", max_rows=5)

        rx = M + left + 0.5
        rw = CONTENT_W - left - 0.5
        self._text(s, rx, y, rw, 0.4, "Items removed from endpoints",
                   size=17, bold=True, color=TEXT_MID)
        rrows = [
            [(i.get("name") or "")[:40], (i.get("hostname") or "")[:22],
             i.get("platform") or "", (i.get("risk_level") or "").capitalize(),
             (i.get("reason") or "")[:26]]
            for i in self.r.rows("remediated_items")
        ]
        if rrows:
            self._table(s, rx, y + 0.55, rw,
                        ["Item", "Endpoint", "Platform", "Risk", "Reason"],
                        rrows, col_w=[3.9, 2.3, 1.8, 1.3, 2.5],
                        max_rows=10, risk_col=3)
        else:
            self._empty_card(s, rx, y + 0.55, rw,
                             "No completed remediations recorded in this window")
        self._footer(s)
        return s

    def slide_agentic(self):
        s = self._slide()
        self._header(s, "Agentic runtime \u2014 what your AI agents did",
                     "Coding-agent sessions observed and governed on managed endpoints",
                     tag="AGENTIC AI", tag_color=AMBER)
        y = 2.5
        cw = (CONTENT_W - 3 * 0.35) / 4
        d = self.r.mapping("agent_decisions")
        measured = self.r.measured("agent_decisions")
        kpi_h = 2.5

        def dec(key):
            return _n(d.get(key, 0)) if measured else NOT_MEASURED

        self._kpi(s, M, y, cw, kpi_h, self.r.num("agent_sessions_total"),
                  "Agent sessions", "last 30 days", TEAL)
        self._kpi(s, M + cw + 0.35, y, cw, kpi_h, dec("allow"),
                  "Actions allowed", "normal developer flow", MINT)
        self._kpi(s, M + 2 * (cw + 0.35), y, cw, kpi_h, dec("ask"),
                  "Actions escalated", "required confirmation", AMBER)
        self._kpi(s, M + 3 * (cw + 0.35), y, cw, kpi_h, dec("block"),
                  "Actions blocked", "policy enforcement", RED)

        y2 = y + kpi_h + 0.45
        h2 = SH - y2 - 0.9
        half = (CONTENT_W - 0.5) / 2
        self._card(s, M, y2, half, h2)
        self._bar_chart(s, M + 0.5, y2 + 0.5, half - 1.0, (h2 - 1.4) / 2,
                        self.r.get("agents_seen"), TEAL,
                        title="Agents in use", max_rows=4)
        self._bar_chart(s, M + 0.5, y2 + 0.9 + (h2 - 1.4) / 2, half - 1.0,
                        (h2 - 1.4) / 2, self.r.get("agent_models"), MINT,
                        title="Models observed", max_rows=4)

        bx = M + half + 0.5
        self._text(s, bx, y2, half, 0.4, "Recent blocked actions",
                   size=17, bold=True, color=TEXT_MID)
        brows = [
            [(b.get("agent") or "")[:16], (b.get("host") or "")[:20],
             (b.get("action") or "")[:14], (b.get("target") or "")[:40]]
            for b in self.r.rows("agent_blocked_examples")
        ]
        if brows:
            self._table(s, bx, y2 + 0.55, half,
                        ["Agent", "Endpoint", "Action", "Target"],
                        brows, col_w=[1.9, 2.3, 1.7, 3.5], max_rows=7)
        else:
            self._empty_card(s, bx, y2 + 0.55, half,
                             "No blocked agent actions in the last 24 hours")
        self._footer(s,
                     "Sessions: 30-day window \u00b7 events: 24-hour window (API limits)")
        return s

    def slide_xsiam(self):
        """Only when an XSIAM correlation was run."""
        x = self.r.xsiam
        if not x:
            return None
        s = self._slide()
        cov = x.get("coverage") or {}
        inc = x.get("incidents") or {}
        self._header(s, "Cortex XSIAM \u2014 where exposure meets detection",
                     "Co-presence on the same estate, not a causal link",
                     tag="CORRELATION", tag_color=TEAL)
        y = 2.5
        cw = (CONTENT_W - 3 * 0.35) / 4
        self._kpi(s, M, y, cw, 2.4, _n(cov.get("on_both")),
                  "Hosts in both systems", "Koi and XSIAM managed", MINT)
        self._kpi(s, M + cw + 0.35, y, cw, 2.4, _n(cov.get("koi_only_count")),
                  "Koi only", "no XSIAM agent seen", AMBER)
        self._kpi(s, M + 2 * (cw + 0.35), y, cw, 2.4,
                  _n(cov.get("xsiam_only_count")),
                  "XSIAM only", "no Koi inventory seen", AMBER)
        self._kpi(s, M + 3 * (cw + 0.35), y, cw, 2.4, _n(inc.get("total")),
                  "XSIAM incidents",
                  "last " + str(x.get("window_days", 30)) + " days", RED)

        y2 = y + 2.85
        half = (CONTENT_W - 0.5) / 2
        self._card(s, M, y2, half, SH - y2 - 0.9)
        self._bar_chart(s, M + 0.5, y2 + 0.5, half - 1.0, SH - y2 - 2.0,
                        inc.get("by_severity") or {}, RED,
                        title="Incidents by severity", max_rows=6)

        bx = M + half + 0.5
        self._text(s, bx, y2, half, 0.4, "Koi-known hosts with incidents",
                   size=17, bold=True, color=TEXT_MID)
        rows = [[(h.get("host") or "")[:38], _n(h.get("incidents"))]
                for h in (inc.get("top_koi_hosts") or [])[:8]]
        if rows:
            self._table(s, bx, y2 + 0.55, half, ["Endpoint", "Incidents"],
                        rows, col_w=[6.4, 2.7], max_rows=8)
        else:
            self._empty_card(s, bx, y2 + 0.55, half,
                             "No XSIAM incident landed on a Koi-known host")
        self._footer(s, "Co-presence on a host is co-presence, not causality")
        return s

    def slide_next_steps(self):
        s = self._slide()
        self._header(s, "Recommended next steps",
                     "Where to go from here, based on what the PoV surfaced",
                     tag="NEXT STEPS", tag_color=TEAL)
        y = 2.6
        recs = self._build_recommendations()[:6]
        cw = (CONTENT_W - 0.5) / 2
        avail = SH - y - 0.9
        rows = (len(recs) + 1) // 2
        card_h = min(2.6, (avail - (rows - 1) * 0.35) / max(rows, 1))
        for i, (title, bodytext, accent) in enumerate(recs):
            col = i % 2
            row = i // 2
            x = M + col * (cw + 0.5)
            ry = y + row * (card_h + 0.35)
            self._card(s, x, ry, cw, card_h)
            bar = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(ry),
                Inches(0.11), Inches(card_h),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.line.fill.background()
            self._no_shadow(bar)
            self._text(s, x + 0.5, ry + 0.32, cw - 0.9, 0.45, title,
                       size=19, bold=True, color=accent)
            self._text(s, x + 0.5, ry + 0.92, cw - 0.9, card_h - 1.2, bodytext,
                       size=14, color=TEXT_MID, line_spacing=1.2)
        self._footer(s)
        return s

    def _build_recommendations(self):
        """Operator text when supplied, otherwise data-driven defaults."""
        written = (self.n.get("recommendations") or "").strip()
        if written:
            out = []
            blocks = [b.strip() for b in written.split("\n") if b.strip()]
            accents = [RED, AMBER, MINT, TEAL, MINT, TEAL]
            for i, block in enumerate(blocks[:6]):
                if ". " in block:
                    head, tail = block.split(". ", 1)
                else:
                    head, tail = "Action " + str(i + 1), block
                out.append((head.strip()[:60], tail.strip(), accents[i % len(accents)]))
            return out

        r = self.r
        out = []
        crit = r.critical
        ungoverned = r.get("ungoverned_high_risk") or 0
        if crit:
            plural = "s" if crit > 1 else ""
            out.append((
                "Remove " + _n(crit) + " critical item" + plural,
                "Critical-risk items were found on managed endpoints. Enable the "
                "Malware Protection guardrail so these are removed automatically "
                "and cannot be reinstalled.",
                RED,
            ))
        if ungoverned:
            out.append((
                "Govern " + _n(ungoverned) + " high-risk items",
                "These items are high or critical risk but no policy currently "
                "governs them. Run an Impact Check, then move them under an "
                "explicit allow or block policy.",
                AMBER,
            ))
        if not r.rows("runtime_policies"):
            out.append((
                "Turn on agent guardrails",
                "No agent runtime policies were active. Enable credential-access "
                "and destructive-command guardrails on the coding agents in use "
                "to cover the agentic attack surface.",
                AMBER,
            ))
        elif r.mapping("agent_decisions").get("block", 0):
            out.append((
                "Extend agent enforcement",
                "Agent guardrails already blocked real actions during the PoV. "
                "Extend the scope from the pilot group to the wider developer "
                "population.",
                MINT,
            ))
        if r.measured("devices_total") and r.raw("devices_total"):
            out.append((
                "Scale beyond the PoV fleet",
                "The PoV covered " + r.num("devices_total") + " endpoints. Extend "
                "deployment through your existing EDR or MDM to the full estate; "
                "no new agent is required.",
                MINT,
            ))
        if (r.measured("blocklist_count") and not r.raw("blocklist_count")
                and (crit or r.high)):
            out.append((
                "Formalise a blocklist",
                "Risky items were found but nothing is explicitly blocklisted. "
                "Codify the decisions made during the PoV so they persist.",
                TEAL,
            ))
        out.append((
            "Enable the Version Update Cooldown",
            "Delaying new versions by a few days defends against the supply-chain "
            "pattern where a trusted package turns malicious in an update.",
            TEAL,
        ))
        if not self.r.xsiam:
            out.append((
                "Connect to Cortex XSIAM",
                "Forward Koi alerts into XSIAM so marketplace risk is correlated "
                "with endpoint and identity telemetry in a single incident view.",
                TEAL,
            ))
        return out

    def slide_appendix(self):
        """Collection notes: warnings and domains never collected."""
        warnings = self.r.raw("warnings", []) or []
        missing = sorted(self.r.missing)
        if not warnings and not missing:
            return None
        s = self._slide()
        self._header(s, "Data collection notes",
                     "What could not be measured during this PoV",
                     tag="APPENDIX", tag_color=TEXT_LO)
        y = 2.5
        if missing:
            self._text(s, M, y, CONTENT_W, 0.4, "Domains not collected",
                       size=17, bold=True, color=TEXT_MID)
            self._text(s, M, y + 0.5, CONTENT_W, 0.4,
                       "\u25b8  " + ", ".join(missing)
                       + "  \u2014  figures for these areas are absent, not zero.",
                       size=13, color=AMBER)
            y += 1.2
        if warnings:
            self._text(s, M, y, CONTENT_W, 0.4, "Collection warnings",
                       size=17, bold=True, color=TEXT_MID)
            for i, w in enumerate(warnings[:10]):
                self._text(s, M, y + 0.5 + i * 0.42, CONTENT_W, 0.38,
                           "\u25b8  " + str(w), size=12, color=TEXT_MID, font=MONO)
        self._footer(s)
        return s

    # ------------------------------------------------------------- build
    def build(self, path: str) -> str:
        self.slide_title()
        self.slide_exec_summary()
        self.slide_narrative()
        self.slide_scorecard()
        self.slide_discovery()
        self.slide_risk()
        self.slide_supply_chain()
        self.slide_supply_chain_items()
        self.slide_threat_intel()
        self.slide_governance()
        self.slide_remediation()
        self.slide_agentic()
        self.slide_xsiam()
        self.slides_findings()
        self.slides_scenarios()
        self.slides_actions()
        self.slide_threat_context()
        self.slide_next_steps()
        self.slide_appendix()
        self.prs.save(path)
        log.info("Deck written to %s", path)
        return path


deck_sections.attach(DeckBuilder)


def build_deck(data: dict, path: str, narrative: dict | None = None) -> str:
    return DeckBuilder(data, narrative).build(path)
