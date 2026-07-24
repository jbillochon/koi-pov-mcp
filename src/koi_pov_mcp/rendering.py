"""
Deliverable rendering from a tenant's aggregated PoV JSON (+ enrichment).

All three formats are pure Python and work on Windows, Linux and macOS with
no system libraries: DOCX (python-docx), PPTX (python-pptx), PDF (reportlab).
WeasyPrint is used instead when it is importable AND reportlab is not, since
it needs GTK on Windows and is therefore not a dependency we can rely on.

Every number rendered here comes from the JSON. Narrative sections
(executive summary, success criteria verdicts, recommendations) are passed
in by the caller and default to visible [[TO BE PROVIDED]] placeholders:
a missing narrative must be seen, not papered over.
"""

from __future__ import annotations

from pathlib import Path

PLACEHOLDER_SUMMARY = "[[TO BE PROVIDED: executive summary]]"
PLACEHOLDER_CRITERIA = "[[TO BE PROVIDED: success criteria scorecard]]"
PLACEHOLDER_RECO = "[[TO BE PROVIDED: recommendations]]"
NO_DATA = "[[TO BE PROVIDED: no data collected for this section]]"


# ---------------------------------------------------------------------- #
# Shared content assembly (single source for all three formats)
# ---------------------------------------------------------------------- #


def _s(v) -> str:
    """Render a value for display. None means not measured, never 0."""
    if v is None:
        return "not measured"
    return str(v)


def _kv_rows(d: dict | None) -> list[list[str]]:
    return [[k, _s(v)] for k, v in (d or {}).items()]


def build_content(data: dict, executive_summary: str, recommendations: str,
                  success_criteria: list[dict] | None) -> dict:
    meta = data.get("meta") or {}
    derived = data.get("derived") or {}
    enr = data.get("enrichment") or {}
    missing = data.get("missing_domains") or []

    top_rows = [
        [_s(i.get("name")), _s(i.get("publisher")), _s(i.get("marketplace")),
         _s(i.get("risk")), _s(i.get("risk_level")), _s(i.get("endpoints"))]
        for i in (data.get("top_risk_items") or [])[:10]
    ]
    action_rows = [
        [_s(i.get("name")), _s(i.get("marketplace")), _s(i.get("risk_level")),
         _s(i.get("endpoints")), ", ".join(map(_s, (i.get("findings") or [])[:3]))]
        for i in (data.get("action_candidates") or [])[:10]
    ]
    remed_rows = [
        [_s(r.get("name")), _s(r.get("hostname")), _s(r.get("platform")),
         _s(r.get("risk_level")), _s(r.get("reason"))]
        for r in (data.get("remediated_items") or [])[:10]
    ]
    policy_rows = [
        [_s(p.get("name")), _s(p.get("action")),
         "enabled" if p.get("enabled") else "disabled", _s(p.get("groups"))]
        for p in (data.get("policies") or [])[:12]
    ]
    blocked_rows = [
        [_s(e.get("agent")), _s(e.get("host")), _s(e.get("action")),
         _s(e.get("target")), _s(e.get("timestamp"))]
        for e in (data.get("agent_blocked_examples") or [])[:8]
    ]

    # TI table: KEV first, then EPSS desc; only enriched facts, dated.
    cves = enr.get("cves") or {}
    ranked = sorted(
        cves.values(),
        key=lambda c: (bool(c.get("kev")), c.get("epss") or 0.0,
                       c.get("cvss_score") or 0.0),
        reverse=True,
    )
    ti_rows = [
        [_s(c.get("id")),
         "YES" + (f" ({c.get('kev_date_added')})" if c.get("kev_date_added") else "")
         if c.get("kev") else "no",
         f"{c['epss']:.2%}" if isinstance(c.get("epss"), float) else "",
         _s(c.get("cvss_score")) if c.get("cvss_score") is not None else "",
         _s(c.get("cvss_severity")) if c.get("cvss_severity") else ""]
        for c in ranked[:10]
    ]
    osv_rows = [
        [pkg, ", ".join(ids[:4])]
        for pkg, ids in list((enr.get("osv") or {}).get("matches", {}).items())[:10]
    ]

    criteria_rows = None
    if success_criteria:
        criteria_rows = [
            [_s(c.get("criterion")), _s(c.get("verdict")), _s(c.get("evidence"))]
            for c in success_criteria
        ]

    return {
        "title": f"Cortex AES Proof of Value: {meta.get('customer_name', 'Customer')}",
        "subtitle": " | ".join(filter(None, [
            f"PoV window: {meta.get('pov_start') or '?'} to {meta.get('pov_end') or '?'}",
            f"Prepared by {meta.get('prepared_by')}" if meta.get("prepared_by") else "",
            f"Generated {meta.get('generated_at', '')}",
        ])),
        "stage": derived.get("stage") or "not determinable (partial collection)",
        "missing_domains": missing,
        "executive_summary": executive_summary.strip() or PLACEHOLDER_SUMMARY,
        "recommendations": recommendations.strip() or PLACEHOLDER_RECO,
        "criteria_rows": criteria_rows,
        "kpis": [
            ("Devices enrolled", _s(data.get("devices_total"))),
            ("Active devices", _s(data.get("devices_active"))),
            ("Items discovered", _s(data.get("items_total"))),
            ("Unique publishers", _s(data.get("unique_publishers"))),
            ("High + critical", _s(derived.get("high_and_critical"))),
            ("Ungoverned high risk", _s(data.get("ungoverned_high_risk"))),
            ("Exposed installs", _s(data.get("exposed_installs"))),
            ("Remediations", _s(data.get("remediations_total"))),
            ("Agent sessions", _s(data.get("agent_sessions_total"))),
        ],
        "devices_by_os": _kv_rows(data.get("devices_by_os")),
        "items_by_view": _kv_rows(data.get("items_by_view")),
        "items_by_marketplace": _kv_rows(data.get("items_by_marketplace")),
        "items_by_risk": _kv_rows(data.get("items_by_risk")),
        "finding_rows": [[_s(f.get("finding")), _s(f.get("count"))]
                         for f in (data.get("finding_frequency") or [])[:12]],
        "top_rows": top_rows,
        "action_rows": action_rows,
        "policy_rows": policy_rows,
        "runtime_policies": [[_s(p.get("name")), _s(p.get("mode")),
                              "enabled" if p.get("enabled") else "disabled",
                              ", ".join(map(_s, p.get("agents") or []))]
                             for p in (data.get("runtime_policies") or [])[:10]],
        "lists": [("Allowlist entries", _s(data.get("allowlist_count"))),
                  ("Blocklist entries", _s(data.get("blocklist_count")))],
        "remed_by_status": _kv_rows(data.get("remediations_by_status")),
        "remed_rows": remed_rows,
        "approvals": _kv_rows(data.get("approvals_by_status")),
        "agent_decisions": _kv_rows(data.get("agent_decisions")),
        "agents_seen": _kv_rows(data.get("agents_seen")),
        "blocked_rows": blocked_rows,
        "alerts_by_severity": _kv_rows(data.get("alerts_by_severity")),
        "ti_rows": ti_rows,
        "osv_rows": osv_rows,
        "ti_fetched_at": enr.get("fetched_at", ""),
        "mitre": enr.get("mitre") or {},
    }


def _mitre_rows(content: dict) -> list[list[str]]:
    rows = []
    for finding, payload in (content.get("mitre") or {}).items():
        techniques = ", ".join(
            f"{t['id']} {t.get('name', '')}".strip()
            for t in payload.get("techniques", [])
        )
        rows.append([finding, techniques])
    return rows[:12]


# ---------------------------------------------------------------------- #
# DOCX
# ---------------------------------------------------------------------- #


def render_docx(content: dict, path: Path) -> None:
    from docx import Document

    doc = Document()

    def table(headers: list[str], rows: list[list[str]]) -> None:
        if not rows:
            doc.add_paragraph(NO_DATA)
            return
        t = doc.add_table(rows=1, cols=len(headers))
        t.style = "Light Grid Accent 1"
        for j, h in enumerate(headers):
            t.rows[0].cells[j].text = h
        for row in rows:
            cells = t.add_row().cells
            for j, v in enumerate(row):
                cells[j].text = _s(v)

    doc.add_heading(content["title"], level=0)
    doc.add_paragraph(content["subtitle"])
    doc.add_paragraph(f"Tenant lifecycle stage: {content['stage']}")
    if content["missing_domains"]:
        doc.add_paragraph(
            "Not collected during this PoV: "
            + ", ".join(content["missing_domains"])
            + ". Figures for these areas are absent, not zero."
        )

    doc.add_heading("Executive summary", level=1)
    doc.add_paragraph(content["executive_summary"])

    doc.add_heading("Success criteria scorecard", level=1)
    if content["criteria_rows"]:
        table(["Criterion", "Verdict", "Evidence"], content["criteria_rows"])
    else:
        doc.add_paragraph(PLACEHOLDER_CRITERIA)

    doc.add_heading("Scope and coverage", level=1)
    table(["Metric", "Value"], [list(k) for k in content["kpis"]])
    doc.add_heading("Devices by OS", level=2)
    table(["OS", "Devices"], content["devices_by_os"])

    doc.add_heading("Discovery", level=1)
    doc.add_heading("Inventory by category", level=2)
    table(["Category", "Items"], content["items_by_view"])
    doc.add_heading("Inventory by marketplace", level=2)
    table(["Marketplace", "Items"], content["items_by_marketplace"])

    doc.add_heading("Risk analysis", level=1)
    table(["Risk level", "Items"], content["items_by_risk"])
    doc.add_heading("Most frequent findings", level=2)
    table(["Finding", "Occurrences"], content["finding_rows"])
    doc.add_heading("Top risk items", level=2)
    table(["Item", "Publisher", "Marketplace", "Risk", "Level", "Endpoints"],
          content["top_rows"])

    doc.add_heading("Threat intelligence", level=1)
    if content["ti_rows"] or content["osv_rows"] or content["mitre"]:
        doc.add_paragraph(
            f"Threat intel as of {content['ti_fetched_at']}. "
            "Priority reads exploitation first: KEV (known exploited), then "
            "EPSS (probability), then CVSS (severity)."
        )
        doc.add_heading("CVEs in scope", level=2)
        table(["CVE", "KEV", "EPSS", "CVSS", "Severity"], content["ti_rows"])
        doc.add_heading("Vulnerable packages (OSV)", level=2)
        table(["Package", "Advisories"], content["osv_rows"])
        doc.add_heading("ATT&CK techniques from findings", level=2)
        table(["Finding", "Techniques"], _mitre_rows(content))
    else:
        doc.add_paragraph("[[TO BE PROVIDED: TI enrichment (run koi_enrich)]]")

    doc.add_heading("Exposure and priority actions", level=1)
    table(["Item", "Marketplace", "Level", "Endpoints", "Findings"],
          content["action_rows"])

    doc.add_heading("Governance", level=1)
    table(["Policy", "Action", "State", "Groups"], content["policy_rows"])
    doc.add_heading("Runtime (agentic) policies", level=2)
    table(["Policy", "Mode", "State", "Agents"], content["runtime_policies"])
    table(["List", "Entries"], [list(k) for k in content["lists"]])

    doc.add_heading("Remediation", level=1)
    table(["Status", "Count"], content["remed_by_status"])
    doc.add_heading("Remediated items (sample)", level=2)
    table(["Item", "Host", "Platform", "Level", "Reason"], content["remed_rows"])

    doc.add_heading("Agentic runtime activity", level=1)
    table(["Decision", "Count"], content["agent_decisions"])
    table(["Agent", "Sessions"], content["agents_seen"])
    doc.add_heading("Blocked actions (sample)", level=2)
    table(["Agent", "Host", "Action", "Target", "Time"], content["blocked_rows"])

    doc.add_heading("Recommendations", level=1)
    doc.add_paragraph(content["recommendations"])

    doc.save(str(path))


# ---------------------------------------------------------------------- #
# PPTX
# ---------------------------------------------------------------------- #


def render_pptx(content: dict, path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    def slide(title: str):
        s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.9))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        return s

    def bullets(s, lines: list[str], size: int = 16, top: float = 1.3):
        tb = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(5.4))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(size)

    def table(s, headers: list[str], rows: list[list[str]], top: float = 1.3):
        if not rows:
            bullets(s, [NO_DATA])
            return
        rows = rows[:9]
        shape = s.shapes.add_table(
            len(rows) + 1, len(headers),
            Inches(0.5), Inches(top), Inches(9), Inches(0.35 * (len(rows) + 1)),
        ).table
        for j, h in enumerate(headers):
            shape.cell(0, j).text = h
        for i, row in enumerate(rows, 1):
            for j, v in enumerate(row):
                shape.cell(i, j).text = _s(v)[:60]

    s = slide(content["title"])
    bullets(s, [content["subtitle"], f"Tenant lifecycle stage: {content['stage']}"], 18)

    bullets(slide("Executive summary"), content["executive_summary"].split("\n"), 18)

    s = slide("Success criteria scorecard")
    if content["criteria_rows"]:
        table(s, ["Criterion", "Verdict", "Evidence"], content["criteria_rows"])
    else:
        bullets(s, [PLACEHOLDER_CRITERIA])

    table(slide("Coverage"), ["Metric", "Value"], [list(k) for k in content["kpis"]])
    table(slide("Discovery: inventory by category"), ["Category", "Items"],
          content["items_by_view"])
    table(slide("Risk distribution"), ["Risk level", "Items"], content["items_by_risk"])
    table(slide("Most frequent findings"), ["Finding", "Occurrences"],
          content["finding_rows"])
    table(slide("Top risk items"),
          ["Item", "Publisher", "Marketplace", "Risk", "Level", "Endpoints"],
          content["top_rows"])

    s = slide("Threat intelligence")
    if content["ti_rows"]:
        bullets(s, [f"As of {content['ti_fetched_at']} - KEV > EPSS > CVSS"], 14)
        table(s, ["CVE", "KEV", "EPSS", "CVSS", "Severity"], content["ti_rows"], top=1.8)
    else:
        bullets(s, ["[[TO BE PROVIDED: TI enrichment]]"])

    s = slide("ATT&CK techniques observed")
    table(s, ["Finding", "Techniques"], _mitre_rows(content))

    table(slide("Exposure: ungoverned high risk"),
          ["Item", "Marketplace", "Level", "Endpoints", "Findings"],
          content["action_rows"])
    table(slide("Governance"), ["Policy", "Action", "State", "Groups"],
          content["policy_rows"])
    table(slide("Remediation"), ["Status", "Count"], content["remed_by_status"])
    table(slide("Agentic runtime"), ["Decision", "Count"], content["agent_decisions"])

    bullets(slide("Recommendations"), content["recommendations"].split("\n"), 18)
    bullets(slide("Next steps"),
            ["[[TO BE PROVIDED: next steps agreed with the customer]]"])

    prs.save(str(path))


# ---------------------------------------------------------------------- #
# PDF: reportlab (pure Python, no system libraries)
# ---------------------------------------------------------------------- #


def _render_pdf_reportlab(content: dict, path: Path) -> None:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
    )

    styles = getSampleStyleSheet()
    body = ParagraphStyle("body", parent=styles["BodyText"], fontSize=9,
                          leading=12, alignment=TA_LEFT)
    cell = ParagraphStyle("cell", parent=body, fontSize=7.5, leading=9.5)
    head_cell = ParagraphStyle("head", parent=cell, textColor=colors.white,
                               fontName="Helvetica-Bold")

    doc = SimpleDocTemplate(
        str(path), pagesize=A4,
        leftMargin=16 * mm, rightMargin=16 * mm,
        topMargin=16 * mm, bottomMargin=16 * mm,
        title=content["title"],
    )
    width = doc.width
    story: list = []

    def h1(text: str) -> None:
        story.append(Spacer(1, 7))
        story.append(Paragraph(text, styles["Heading1"]))

    def h2(text: str) -> None:
        story.append(Spacer(1, 4))
        story.append(Paragraph(text, styles["Heading2"]))

    def para(text: str) -> None:
        for chunk in str(text).split("\n"):
            if chunk.strip():
                story.append(Paragraph(chunk, body))
                story.append(Spacer(1, 3))

    def table(headers: list[str], rows: list[list[str]]) -> None:
        if not rows:
            para(NO_DATA)
            return
        data = [[Paragraph(h, head_cell) for h in headers]]
        for row in rows:
            data.append([Paragraph(_s(v), cell) for v in row])
        col = width / len(headers)
        t = Table(data, colWidths=[col] * len(headers), repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#31445c")),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#9aa5b1")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [colors.white, colors.HexColor("#f2f4f7")]),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]))
        story.append(t)
        story.append(Spacer(1, 6))

    story.append(Paragraph(content["title"], styles["Title"]))
    para(content["subtitle"])
    para(f"Tenant lifecycle stage: {content['stage']}")
    if content["missing_domains"]:
        para("Not collected during this PoV: "
             + ", ".join(content["missing_domains"])
             + ". Figures for these areas are absent, not zero.")

    h1("Executive summary")
    para(content["executive_summary"])

    h1("Success criteria scorecard")
    if content["criteria_rows"]:
        table(["Criterion", "Verdict", "Evidence"], content["criteria_rows"])
    else:
        para(PLACEHOLDER_CRITERIA)

    h1("Scope and coverage")
    table(["Metric", "Value"], [list(k) for k in content["kpis"]])
    h2("Devices by OS")
    table(["OS", "Devices"], content["devices_by_os"])

    story.append(PageBreak())
    h1("Discovery")
    h2("Inventory by category")
    table(["Category", "Items"], content["items_by_view"])
    h2("Inventory by marketplace")
    table(["Marketplace", "Items"], content["items_by_marketplace"])

    h1("Risk analysis")
    table(["Risk level", "Items"], content["items_by_risk"])
    h2("Most frequent findings")
    table(["Finding", "Occurrences"], content["finding_rows"])
    h2("Top risk items")
    table(["Item", "Publisher", "Marketplace", "Risk", "Level", "Endpoints"],
          content["top_rows"])

    story.append(PageBreak())
    h1("Threat intelligence")
    if content["ti_rows"] or content["osv_rows"] or content["mitre"]:
        para(f"Threat intel as of {content['ti_fetched_at']}. Priority reads "
             "exploitation first: KEV (known exploited), then EPSS "
             "(probability), then CVSS (severity).")
        h2("CVEs in scope")
        table(["CVE", "KEV", "EPSS", "CVSS", "Severity"], content["ti_rows"])
        h2("Vulnerable packages (OSV)")
        table(["Package", "Advisories"], content["osv_rows"])
        h2("ATT&CK techniques from findings")
        table(["Finding", "Techniques"], _mitre_rows(content))
    else:
        para("[[TO BE PROVIDED: TI enrichment (run koi_enrich)]]")

    h1("Exposure and priority actions")
    table(["Item", "Marketplace", "Level", "Endpoints", "Findings"],
          content["action_rows"])

    story.append(PageBreak())
    h1("Governance")
    table(["Policy", "Action", "State", "Groups"], content["policy_rows"])
    h2("Runtime (agentic) policies")
    table(["Policy", "Mode", "State", "Agents"], content["runtime_policies"])
    table(["List", "Entries"], [list(k) for k in content["lists"]])

    h1("Remediation")
    table(["Status", "Count"], content["remed_by_status"])
    h2("Remediated items (sample)")
    table(["Item", "Host", "Platform", "Level", "Reason"], content["remed_rows"])

    h1("Agentic runtime activity")
    table(["Decision", "Count"], content["agent_decisions"])
    table(["Agent", "Sessions"], content["agents_seen"])
    h2("Blocked actions (sample)")
    table(["Agent", "Host", "Action", "Target", "Time"], content["blocked_rows"])

    h1("Recommendations")
    para(content["recommendations"])

    doc.build(story)


def _render_pdf_weasyprint(content: dict, path: Path) -> None:
    """Fallback engine. Needs GTK on Windows, hence not the default."""
    from weasyprint import HTML

    def html_table(headers, rows):
        if not rows:
            return f"<p><em>{NO_DATA}</em></p>"
        head = "".join(f"<th>{h}</th>" for h in headers)
        body = "".join(
            "<tr>" + "".join(f"<td>{_s(v)}</td>" for v in row) + "</tr>"
            for row in rows
        )
        return f"<table><tr>{head}</tr>{body}</table>"

    sections = [
        ("Executive summary", f"<p>{content['executive_summary']}</p>"),
        ("Success criteria scorecard",
         html_table(["Criterion", "Verdict", "Evidence"], content["criteria_rows"] or [])),
        ("Coverage", html_table(["Metric", "Value"], [list(k) for k in content["kpis"]])),
        ("Risk distribution", html_table(["Risk level", "Items"], content["items_by_risk"])),
        ("Top risk items",
         html_table(["Item", "Publisher", "Marketplace", "Risk", "Level", "Endpoints"],
                    content["top_rows"])),
        ("Threat intelligence", html_table(["CVE", "KEV", "EPSS", "CVSS", "Severity"],
                                           content["ti_rows"])),
        ("Exposure", html_table(["Item", "Marketplace", "Level", "Endpoints", "Findings"],
                                content["action_rows"])),
        ("Governance", html_table(["Policy", "Action", "State", "Groups"],
                                  content["policy_rows"])),
        ("Remediation", html_table(["Status", "Count"], content["remed_by_status"])),
        ("Agentic runtime", html_table(["Decision", "Count"], content["agent_decisions"])),
        ("Recommendations", f"<p>{content['recommendations']}</p>"),
    ]
    body = "".join(f"<h2>{t}</h2>{h}" for t, h in sections)
    html = f"""<html><head><meta charset="utf-8"><style>
    body {{ font-family: Helvetica, Arial, sans-serif; font-size: 11px; }}
    h1 {{ font-size: 20px; }} h2 {{ font-size: 14px; margin-top: 18px; }}
    table {{ border-collapse: collapse; width: 100%; }}
    th, td {{ border: 1px solid #999; padding: 3px 6px; text-align: left; }}
    th {{ background: #eee; }}
    </style></head><body>
    <h1>{content['title']}</h1><p>{content['subtitle']}</p>{body}
    </body></html>"""
    HTML(string=html).write_pdf(str(path))


def render_pdf(content: dict, path: Path) -> None:
    """reportlab first (no system libraries); WeasyPrint only as a fallback."""
    try:
        _render_pdf_reportlab(content, path)
        return
    except ImportError:
        pass
    _render_pdf_weasyprint(content, path)


def render(data: dict, out_dir: Path, formats: list[str],
           executive_summary: str = "", recommendations: str = "",
           success_criteria: list[dict] | None = None) -> dict:
    """Render the requested formats. Returns {'produced': {...}, 'skipped': {...}}."""
    content = build_content(data, executive_summary, recommendations, success_criteria)
    out_dir.mkdir(parents=True, exist_ok=True)
    produced: dict[str, str] = {}
    skipped: dict[str, str] = {}

    for fmt in formats:
        fmt = fmt.lower().strip()
        try:
            if fmt == "docx":
                p = out_dir / "report.docx"
                render_docx(content, p)
                produced["docx"] = str(p)
            elif fmt == "pptx":
                p = out_dir / "deck.pptx"
                render_pptx(content, p)
                produced["pptx"] = str(p)
            elif fmt == "pdf":
                p = out_dir / "report.pdf"
                render_pdf(content, p)
                produced["pdf"] = str(p)
            else:
                skipped[fmt] = "unknown format (docx, pptx, pdf)"
        except ImportError as exc:
            skipped[fmt] = f"missing dependency: {exc}"
        except Exception as exc:  # noqa: BLE001 - one format must not sink the others
            skipped[fmt] = f"{type(exc).__name__}: {exc}"

    return {"produced": produced, "skipped": skipped}
